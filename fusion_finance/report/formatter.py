from __future__ import annotations

import json
import logging
from datetime import datetime
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

TEMPLATES_DIR = Path(__file__).parent / "templates"

TEMPLATE_MAP = {
    "valuation": "valuation.html",
    "pitchbook": "pitchbook.html",
    "research": "research.html",
    "board_material": "board_material.html",
}

FORMAT_HTML = "html"
FORMAT_PDF = "pdf"
FORMAT_PPTX = "pptx"
FORMAT_XLSX = "xlsx"
FORMAT_JSON = "json"
FORMAT_MD = "markdown"

SUPPORTED_FORMATS = (FORMAT_HTML, FORMAT_PDF, FORMAT_PPTX, FORMAT_XLSX, FORMAT_JSON, FORMAT_MD)


class ReportFormatter:
    def __init__(self):
        self._jinja_env = None
        self._weasyprint = None
        self._pptx = None
        self._openpyxl = None
        self._detect_deps()

    def _detect_deps(self):
        try:
            import weasyprint

            self._weasyprint = weasyprint
            logger.debug("WeasyPrint available for PDF export")
        except ImportError:
            logger.debug("WeasyPrint not available, PDF export disabled")
        try:
            from pptx import Presentation

            self._pptx = Presentation
            logger.debug("python-pptx available for PPTX export")
        except ImportError:
            logger.debug("python-pptx not available, PPTX export disabled")
        try:
            import openpyxl

            self._openpyxl = openpyxl
            logger.debug("openpyxl available for XLSX export")
        except ImportError:
            logger.debug("openpyxl not available, XLSX export disabled")

    def _get_jinja(self):
        if self._jinja_env is not None:
            return self._jinja_env
        try:
            from jinja2 import Environment, FileSystemLoader

            self._jinja_env = Environment(
                loader=FileSystemLoader(str(TEMPLATES_DIR)),
                autoescape=True,
            )
            self._jinja_env.filters["safe"] = lambda v: v
            logger.debug("Jinja2 environment initialized")
        except ImportError:
            logger.warning("Jinja2 not available, HTML template rendering disabled")
            self._jinja_env = False
        return self._jinja_env

    def render_html(self, template_name: str, data: dict[str, Any]) -> str:
        env = self._get_jinja()
        if not env:
            return self._fallback_html(template_name, data)
        filename = TEMPLATE_MAP.get(template_name)
        if not filename:
            logger.error("Unknown template: %s", template_name)
            return self._fallback_html(template_name, data)
        try:
            tmpl = env.get_template(filename)
            tmpl.globals["safe"] = lambda v: v
            html = tmpl.render(**data)
            logger.info("Rendered HTML template: %s", template_name)
            return html
        except Exception as e:
            logger.error("Template render failed for %s: %s", template_name, e)
            return self._fallback_html(template_name, data)

    def _fallback_html(self, template_name: str, data: dict[str, Any]) -> str:
        company = data.get("company", "Unknown")
        body = data.get("body", data.get("content", json.dumps(data, ensure_ascii=False, default=str)))
        date_str = data.get("date", datetime.now().strftime("%Y-%m-%d"))
        return f"""<!DOCTYPE html>
<html lang="zh-CN"><head><meta charset="UTF-8"><title>{company} - {template_name}</title>
<style>body{{font-family:sans-serif;background:#0f0f1a;color:#e0e0e0;padding:40px;max-width:900px;margin:0 auto;}}</style>
</head><body><h1>{company} - {template_name}</h1><p>{date_str}</p><div style="white-space:pre-wrap">{body}</div></body></html>"""

    def export(
        self,
        content: str,
        fmt: str,
        output_path: str = "",
        template_name: str = "",
        template_data: dict[str, Any] | None = None,
    ) -> str:
        fmt = fmt.lower()
        if fmt not in SUPPORTED_FORMATS:
            raise ValueError(f"Unsupported format: {fmt}. Supported: {SUPPORTED_FORMATS}")
        if not output_path:
            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
            ext = fmt if fmt != FORMAT_MD else "md"
            output_path = f"report_{ts}.{ext}"
        path = Path(output_path).expanduser()
        path.parent.mkdir(parents=True, exist_ok=True)
        if fmt == FORMAT_HTML:
            html = content
            if template_name and template_data:
                html = self.render_html(template_name, template_data)
            path.write_text(html, encoding="utf-8")
            logger.info("Exported HTML: %s", path)
        elif fmt == FORMAT_PDF:
            html = content
            if template_name and template_data:
                html = self.render_html(template_name, template_data)
            self._export_pdf(html, path)
        elif fmt == FORMAT_PPTX:
            self._export_pptx(content, path, template_data)
        elif fmt == FORMAT_XLSX:
            self._export_xlsx(content, path, template_data)
        elif fmt == FORMAT_JSON:
            data = template_data or {"content": content}
            path.write_text(json.dumps(data, ensure_ascii=False, indent=2, default=str), encoding="utf-8")
            logger.info("Exported JSON: %s", path)
        elif fmt == FORMAT_MD:
            path.write_text(content, encoding="utf-8")
            logger.info("Exported Markdown: %s", path)
        return str(path)

    def _export_pdf(self, html: str, path: Path):
        if not self._weasyprint:
            logger.warning("WeasyPrint unavailable, falling back to HTML")
            html_path = path.with_suffix(".html")
            html_path.write_text(html, encoding="utf-8")
            logger.info("PDF unavailable, saved HTML instead: %s", html_path)
            return
        try:
            doc = self._weasyprint.HTML(string=html)
            doc.write_pdf(str(path))
            logger.info("Exported PDF: %s", path)
        except Exception as e:
            logger.error("PDF export failed: %s", e)
            html_path = path.with_suffix(".html")
            html_path.write_text(html, encoding="utf-8")

    def _export_pptx(self, content: str, path: Path, data: dict[str, Any] | None = None):
        if not self._pptx:
            logger.warning("python-pptx unavailable, falling back to text")
            path.with_suffix(".txt").write_text(content, encoding="utf-8")
            return
        try:
            prs = self._pptx()
            title = data.get("company", "Report") if data else "Report"
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            slide.placeholders[1].text = data.get("date", datetime.now().strftime("%Y-%m-%d")) if data else ""
            content_layout = prs.slide_layouts[1]
            for chunk in _split_content(content, max_chars=800):
                cs = prs.slides.add_slide(content_layout)
                cs.shapes.title.text = title
                cs.placeholders[1].text = chunk
            prs.save(str(path))
            logger.info("Exported PPTX: %s", path)
        except Exception as e:
            logger.error("PPTX export failed: %s", e)
            path.with_suffix(".txt").write_text(content, encoding="utf-8")

    def _export_xlsx(self, content: str, path: Path, data: dict[str, Any] | None = None):
        if not self._openpyxl:
            logger.warning("openpyxl unavailable, falling back to CSV")
            csv_path = path.with_suffix(".csv")
            csv_path.write_text(content, encoding="utf-8")
            return
        try:
            wb = self._openpyxl.Workbook()
            ws = wb.active
            ws.title = "Report"
            if data and "dcf_summary" in data:
                ws.append(["指标", "数值"])
                for item in data["dcf_summary"]:
                    ws.append([item.get("label", ""), item.get("value", "")])
            elif data and "financial_table" in data:
                ft = data["financial_table"]
                ws.append(ft.get("columns", []))
                for row in ft.get("rows", []):
                    ws.append(row)
            else:
                for line in content.split("\n"):
                    ws.append([line])
            wb.save(str(path))
            logger.info("Exported XLSX: %s", path)
        except Exception as e:
            logger.error("XLSX export failed: %s", e)
            csv_path = path.with_suffix(".csv")
            csv_path.write_text(content, encoding="utf-8")


def _split_content(text: str, max_chars: int = 800) -> list[str]:
    chunks = []
    current = ""
    for line in text.split("\n"):
        if len(current) + len(line) + 1 > max_chars and current:
            chunks.append(current)
            current = line
        else:
            current = current + "\n" + line if current else line
    if current:
        chunks.append(current)
    return chunks
